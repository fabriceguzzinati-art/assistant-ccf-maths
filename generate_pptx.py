from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path
import textwrap

outdir = Path("output")
outdir.mkdir(exist_ok=True)

logo_path = Path("Logo-Lino.jpg")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

ORANGE = RGBColor(0xF3, 0x6C, 0x21)
BLUE = RGBColor(0x00, 0x4A, 0x7C)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(255, 255, 255)
MIDGRAY = RGBColor(0xEC, 0xEC, 0xEC)
SOFTBLUE = RGBColor(0xE8, 0xF1, 0xF8)
SOFTORANGE = RGBColor(0xFF, 0xEF, 0xE3)

def add_bg(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bottom_bar(slide):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(6.98), Inches(12.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(3), Inches(0.22))
    p = tx.text_frame.paragraphs[0]
    p.text = "atelier IA 2026"
    p.font.name = "Open Sans"
    p.font.size = Pt(11)
    p.font.color.rgb = BLUE

    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(11.72), Inches(6.7), height=Inches(0.44))

def add_title(slide, title, subtitle=None, color=BLUE):
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.85))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Montserrat"
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = color
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Open Sans"
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK

def text_box(slide, left, top, width, height, text, font_size=20, color=DARK, bold=False, align=PP_ALIGN.CENTER):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Open Sans"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx

def card(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line if line else fill
    return shp

# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
card(s, 0.55, 0.55, 12.2, 5.7, WHITE, GRAY)
text_box(s, 1.05, 1.15, 11.25, 0.95, "L’enseignement explicite augmenté par l’IA", 31, BLUE, True)
text_box(s, 1.35, 2.35, 10.6, 0.7, "Comment gagner du temps tout en renforçant la différenciation pédagogique ?", 18, ORANGE)
card(s, 1.2, 3.45, 10.9, 0.75, SOFTBLUE, SOFTBLUE)
text_box(s, 1.45, 3.68, 10.4, 0.25, "Atelier pour enseignants de lycée professionnel", 15, BLUE, True)
if logo_path.exists():
    s.shapes.add_picture(str(logo_path), Inches(5.6), Inches(4.85), height=Inches(1.0))
add_bottom_bar(s)

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Question d’accroche")
text_box(s, 1.0, 1.35, 11.3, 0.65, "Combien de temps vous faut-il pour préparer une séance complète ?", 22, BLUE, True)

for i, (label, col) in enumerate([("1h", ORANGE), ("2h", BLUE), ("Parfois davantage", DARK)]):
    x = 0.95 + i * 4.1
    card(s, x, 2.2, 3.7, 1.7, WHITE, GRAY)
    text_box(s, x, 2.42, 3.7, 0.3, "⏱️", 28, col)
    text_box(s, x + 0.08, 2.87, 3.55, 0.42, label, 19, col, True)

card(s, 1.2, 4.55, 10.95, 0.78, MIDGRAY, MIDGRAY)
text_box(s, 1.45, 4.79, 10.45, 0.22, "Et si nous pouvions réduire ce temps tout en améliorant la différenciation ?", 18, DARK)
add_bottom_bar(s)

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Objectif de l’atelier")
text_box(s, 0.95, 1.25, 11.4, 0.5, "Aujourd’hui vous allez :", 24, BLUE, True)

steps = [
    ("1", "Construire une activité", ORANGE),
    ("2", "La différencier", BLUE),
    ("3", "Créer une évaluation", RGBColor(0x8F, 0x8F, 0x8F)),
    ("4", "Repartir avec une ressource prête à utiliser", ORANGE),
]
for i, (n, txt, col) in enumerate(steps):
    x = 0.6 + i * 3.1
    card(s, x, 2.15, 2.8, 2.25, WHITE, GRAY)
    card(s, x + 1.0, 2.25, 0.65, 0.65, col, col)
    text_box(s, x + 1.0, 2.38, 0.65, 0.2, n, 16, WHITE, True)
    text_box(s, x + 0.12, 3.0, 2.55, 0.75, txt, 14, DARK)

add_bottom_bar(s)

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Les 3 usages les plus rentables")

for x, title, body, col in [
    (0.95, "Préparer", "Séances, consignes, supports", ORANGE),
    (4.45, "Différencier", "Aides, variantes, niveaux", BLUE),
    (7.95, "Évaluer", "Questions, corrigés, barèmes", RGBColor(0x8F, 0x8F, 0x8F)),
]:
    card(s, x, 2.0, 3.15, 2.7, col, col)
    text_box(s, x, 2.4, 3.15, 0.45, title, 22, WHITE, True)
    text_box(s, x + 0.15, 3.0, 2.85, 0.65, body, 14, WHITE)

add_bottom_bar(s)

# Slide 5
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BLUE)
text_box(s, 0.75, 0.45, 12, 0.55, "Démonstration WOW n°1", 24, WHITE, True)
card(s, 0.75, 1.15, 11.85, 5.55, RGBColor(0x08, 0x3A, 0x61), RGBColor(0x08, 0x3A, 0x61))
text_box(s, 1.0, 1.45, 11.2, 0.45, "Prompt à copier-coller", 16, ORANGE, True)

prompt1 = (
    "Tu es enseignant en lycée professionnel. Construis une séance de 55 minutes selon les principes "
    "de l’enseignement explicite. Niveau : [classe] Thème : [thème] "
    "Je veux : objectif / modelage / pratique guidée / pratique autonome / évaluation de sortie."
)
tx = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.2), Inches(3.8))
tf = tx.text_frame
tf.word_wrap = True
for i, line in enumerate(textwrap.wrap(prompt1, 60)):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = line
    r.font.name = "Open Sans"
    r.font.size = Pt(18)
    r.font.color.rgb = WHITE

card(s, 1.0, 6.0, 4.0, 0.38, ORANGE, ORANGE)
text_box(s, 1.0, 6.03, 4.0, 0.16, "Effet : texte qui apparaît progressivement", 12, WHITE)
add_bottom_bar(s)

# Slide 6
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Démonstration WOW n°2")

levels = [
    ("Élève fragile", "Aide renforcée", ORANGE),
    ("Élève standard", "Chemin classique", BLUE),
    ("Élève avancé", "Défi supplémentaire", RGBColor(0x8F, 0x8F, 0x8F)),
]
for i, (h, txt, col) in enumerate(levels):
    x = 0.72 + i * 4.18
    card(s, x, 2.0, 3.95, 2.7, WHITE, GRAY)
    card(s, x, 2.0, 3.95, 0.5, col, col)
    text_box(s, x, 2.08, 3.95, 0.22, h, 16, WHITE, True)
    text_box(s, x + 0.12, 2.85, 3.7, 0.5, txt, 17, DARK, True)
    text_box(s, x + 0.12, 3.45, 3.7, 0.35, "Adapter sans changer l’objectif", 13, DARK)

text_box(s, 1.0, 5.1, 11.4, 0.55, "Reprends cette séance et adapte-la aux trois niveaux. Conserve le même objectif.", 18, DARK)
add_bottom_bar(s)

# Slide 7
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Démonstration WOW n°3")

card(s, 0.95, 1.95, 4.2, 2.95, WHITE, GRAY)
text_box(s, 0.95, 2.35, 4.2, 0.45, "📝", 34, ORANGE)
text_box(s, 1.0, 2.95, 4.1, 0.55, "Créer une évaluation de 15 minutes", 18, DARK, True)

card(s, 5.35, 1.95, 6.95, 2.95, MIDGRAY, MIDGRAY)
text_box(s, 5.7, 2.5, 6.3, 0.45, "Sujet", 22, BLUE, True)
text_box(s, 5.7, 3.05, 6.3, 0.45, "Corrigé", 22, BLUE, True)
text_box(s, 5.7, 3.6, 6.3, 0.45, "Barème", 22, BLUE, True)
add_bottom_bar(s)

# Slide 8
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Activité des collègues")
card(s, 0.95, 1.7, 11.35, 4.05, WHITE, GRAY)
text_box(s, 1.25, 2.2, 10.8, 0.55, "Prenez un cours que vous devez faire dans les 15 prochains jours.", 22, BLUE, True)
text_box(s, 1.25, 2.92, 10.8, 0.55, "Objectif : créer une activité, une différenciation, une évaluation.", 19, DARK)
card(s, 1.25, 3.72, 10.8, 1.2, MIDGRAY, MIDGRAY)
text_box(s, 1.5, 4.05, 10.3, 0.35, "Zone de travail libre / brouillon / prise de notes", 16, DARK)
add_bottom_bar(s)

# Slide 9
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Mutualisation")
card(s, 1.0, 2.0, 11.3, 3.5, WHITE, GRAY)
text_box(s, 1.35, 2.5, 10.6, 0.6, "Partageons vos productions :", 24, BLUE, True)
text_box(s, 1.35, 3.15, 10.6, 0.55, "ce qui vous a surpris, ce que vous allez réutiliser.", 20, DARK)
card(s, 1.35, 4.0, 10.6, 0.85, SOFTBLUE, SOFTBLUE)
text_box(s, 1.55, 4.28, 10.2, 0.2, "Espace pour captures ou retours oraux", 15, BLUE)
add_bottom_bar(s)

# Slide 10
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Bonus final")
card(s, 0.95, 2.0, 4.5, 2.8, WHITE, GRAY)
text_box(s, 0.95, 2.58, 4.5, 0.45, "📄 Document", 26, BLUE, True)
text_box(s, 5.7, 2.95, 1.0, 0.45, "➜", 32, ORANGE, True)
card(s, 6.7, 2.0, 5.25, 2.8, WHITE, GRAY)
text_box(s, 6.7, 2.58, 5.25, 0.45, "✨ Activité pédagogique", 24, ORANGE, True)
text_box(s, 1.15, 5.1, 11.0, 0.55, "Transforme ce document en activité pédagogique.", 22, DARK)
add_bottom_bar(s)

# Slide 11
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BLUE)
text_box(s, 0.85, 1.95, 11.6, 1.25, "L’IA ne remplace pas l’enseignant.\nElle lui rend du temps.", 30, WHITE, True)
card(s, 3.0, 4.6, 7.3, 0.55, ORANGE, ORANGE)
text_box(s, 3.15, 4.76, 7.0, 0.18, "Le temps gagné devient du temps pédagogique.", 16, WHITE, True)
add_bottom_bar(s)

# Slide 12
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title(s, "Ressources & QR codes")
text_box(s, 1.0, 1.55, 11.3, 0.45, "QR vers fiche participant + QR vers ressources IA éducatives", 18, DARK)

for x, label in [(2.35, "Fiche participant"), (7.65, "Ressources IA")]:
    card(s, x, 2.3, 2.9, 2.9, WHITE, GRAY)
    card(s, x + 0.4, 2.65, 2.1, 2.1, MIDGRAY, MIDGRAY)
    text_box(s, x + 0.4, 3.15, 2.1, 0.5, "QR CODE", 18, BLUE, True)
    text_box(s, x, 5.35, 2.9, 0.35, label, 14, DARK, True)

add_bottom_bar(s)

prs.save(str(outdir / "atelier_enseignement_explicite_IA_2026_v4.pptx"))
print("output/atelier_enseignement_explicite_IA_2026_v4.pptx")